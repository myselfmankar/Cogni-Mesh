import os
import fitz  # PyMuPDF
import docx
import pandas as pd
# Removed whisper - no audio/video processing needed
import warnings
import logging
from typing import Dict, Any, Optional
from pptx import Presentation

# Suppress warnings
warnings.filterwarnings("ignore")

class ContentExtractor:
    def __init__(self, whisper_model_size: str = "base"):
        # Keeping parameter for compatibility but not using it
        pass

    def extract(self, file_path: str) -> Optional[Dict[str, Any]]:
        """
        Detects file type and extracts text.
        Returns a dictionary with keys: 'text', 'metadata'.
        Returns None if file type is not supported or extraction fails.
        """
        if not os.path.exists(file_path):
            logging.warning(f"File not found: {file_path}")
            return None

        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.pdf':
                return self._extract_pdf(file_path)
            elif ext == '.docx':
                return self._extract_docx(file_path)
            elif ext == '.pptx':
                return self._extract_pptx(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self._extract_excel(file_path)
            elif ext in ['.mp4', '.mp3', '.wav', '.m4a', '.mov', '.avi', '.mkv']:
                logging.info(f"Skipping audio/video file (Whisper not installed): {file_path}")
                return None
            elif ext == '.txt':
                return self._extract_txt(file_path)
            elif ext in ['.md', '.markdown']:
                return self._extract_markdown(file_path)
            else:
                logging.info(f"Skipping unsupported file type: {ext} for {file_path}")
                return None
        except Exception as e:
            logging.error(f"Error extracting {file_path}: {str(e)}")
            return None

    def _extract_pdf(self, file_path: str) -> Dict[str, Any]:
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        return {
            "text": text,
            "metadata": {"file_type": "pdf", "page_count": len(doc)}
        }

    def _extract_docx(self, file_path: str) -> Dict[str, Any]:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return {
            "text": text,
            "metadata": {"file_type": "docx"}
        }

    def _extract_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = ""
        slide_count = 0
        
        for slide in prs.slides:
            slide_count += 1
            text += f"\n--- Slide {slide_count} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return {
            "text": text,
            "metadata": {"file_type": "pptx", "slide_count": slide_count}
        }

    def _extract_excel(self, file_path: str) -> Dict[str, Any]:
        # Read all sheets
        xls = pd.ExcelFile(file_path)
        text = ""
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
        return {
            "text": text,
            "metadata": {"file_type": "xlsx", "sheets": ", ".join(xls.sheet_names)}
        }


    def _extract_txt(self, file_path: str) -> Dict[str, Any]:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return {
            "text": text,
            "metadata": {"file_type": "txt"}
        }

    def _extract_markdown(self, file_path: str) -> Dict[str, Any]:
        """
        Extract Markdown files while preserving structure.
        Markdown is ideal for RAG as it maintains semantic structure (headers, lists, code blocks).
        """
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        # Count headers for metadata (useful for understanding document structure)
        header_count = text.count('\n#')
        
        return {
            "text": text,
            "metadata": {
                "file_type": "markdown",
                "header_count": header_count,
                "has_code_blocks": "```" in text
            }
        }
